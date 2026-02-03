import streamlit as st
import os
import json
import io
import re
import time
import requests
import base64
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document

# --- 1. CONFIGURATION ---
st.set_page_config(page_title="AI Presentation Architect", layout="wide")

# --- 2. AI AGENTS ---


def get_gemini_response(api_key, prompt, temp=0.3):
    os.environ["GOOGLE_API_KEY"] = api_key
    # FIX: Use the specific version ID to avoid "404 Not Found"
    llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash-001", temperature=temp)
    return llm.invoke(prompt).content


def generate_design_theme(api_key, topic, audience):
    prompt = f"""
    Role: Senior Art Director. Task: Create a custom design system.
    CONTEXT: Topic: "{topic}", Audience: "{audience}"
    OUTPUT FORMAT (JSON ONLY):
    {{
        "theme_name": "Name",
        "bg_hex": "#FFFFFF", "text_hex": "#000000", "accent_hex": "#0000FF", 
        "chart_palette": ["#0000FF", "#00AA00", "#FF0000", "#FFAA00"]
    }}
    """
    response = get_gemini_response(api_key, prompt, temp=0.7)
    try:
        return json.loads(response.replace("```json", "").replace("```", "").strip())
    except:
        return {"theme_name": "Default", "bg_hex": "#FFFFFF", "text_hex": "#000000", "accent_hex": "#0000FF", "chart_palette": ["#0000FF", "#FF0000"]}

# --- 3. VISUAL ENGINE (DIRECT API) ---


def generate_real_image(prompt, theme_name, google_key):
    """
    Direct REST API Call - Bypasses the Python Library entirely.
    """
    if not google_key:
        return None

    full_prompt = f"{prompt}, {theme_name} style, professional presentation graphic, minimalist, high quality, 4k"

    # RAW API ENDPOINT (You confirmed this works!)
    url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-3.0-generate-001:predict?key={google_key}"

    headers = {"Content-Type": "application/json"}
    payload = {
        "instances": [{"prompt": full_prompt}],
        "parameters": {"sampleCount": 1}
    }

    try:
        response = requests.post(url, headers=headers, json=payload, timeout=8)

        if response.status_code == 200:
            result = response.json()
            if 'predictions' in result:
                b64_data = result['predictions'][0]['bytesBase64Encoded']
                return io.BytesIO(base64.b64decode(b64_data))
            return None
        else:
            # If Google fails, silently return None so we fallback to text
            return None

    except Exception as e:
        return None

# --- 4. HELPERS ---


def hex_to_rgb(hex_code):
    hex_code = hex_code.lstrip('#')
    return tuple(int(hex_code[i:i+2], 16) for i in (0, 2, 4))


def apply_markdown_to_paragraph(paragraph, text, text_rgb):
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        run = paragraph.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part
            run.font.bold = False
        run.font.color.rgb = text_rgb
        run.font.size = Pt(20)


def markdown_to_html(text):
    return re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)


def extract_text_from_file(uploaded_file):
    file_type = uploaded_file.name.split('.')[-1].lower()
    text = ""
    try:
        if file_type == 'pdf':
            with open("temp.pdf", "wb") as f:
                f.write(uploaded_file.getbuffer())
            loader = PyPDFLoader("temp.pdf")
            text = "\n".join([p.page_content for p in loader.load()])
        elif file_type == 'docx':
            doc = Document(uploaded_file)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif file_type in ['xlsx', 'xls', 'csv']:
            df = pd.read_excel(
                uploaded_file) if 'xls' in file_type else pd.read_csv(uploaded_file)
            text = df.to_string()
        elif file_type == 'pptx':
            prs = Presentation(uploaded_file)
            text = "\n".join(
                [s.text for s in prs.slides for s in s.shapes if hasattr(s, "text")])
        elif file_type == 'txt':
            text = str(uploaded_file.read(), "utf-8")
    except Exception as e:
        return f"Error: {str(e)}"
    return f"\n\n--- SOURCE: {uploaded_file.name} ---\n{text}"


def create_chart_image(chart_info, theme_data):
    cats = chart_info.get("categories", [])
    vals = chart_info.get("values", [])
    c_type = chart_info.get("type", "BAR")
    if not cats or not vals:
        return None

    colors = theme_data["chart_palette"]
    bg_color = theme_data["bg_hex"]
    text_color = theme_data["text_hex"]

    plt.clf()
    sns.set_theme(style="white")
    fig, ax = plt.subplots(figsize=(6, 4))
    fig.patch.set_facecolor(bg_color)
    ax.set_facecolor(bg_color)
    for spine in ax.spines.values():
        spine.set_color(text_color)
    ax.tick_params(colors=text_color)

    try:
        if c_type == "BAR":
            sns.barplot(x=vals, y=cats,
                        palette=colors[:len(cats)], orient='h', ax=ax)
            sns.despine(left=True, bottom=True)
            ax.set(xticks=[])
            for i, v in enumerate(vals):
                ax.text(v, i, f" {v}", color=text_color,
                        va='center', fontweight='bold')
        elif c_type == "LINE":
            sns.lineplot(x=cats, y=vals, marker="o",
                         linewidth=3, color=colors[0], ax=ax)
            sns.despine(left=True)
            plt.grid(axis='y', linestyle='--', alpha=0.3, color=text_color)
        elif c_type == "PIE":
            plt.pie(vals, labels=cats, colors=colors, autopct='%1.1f%%', startangle=90,
                    textprops={'color': text_color}, wedgeprops=dict(width=0.5))

        plt.title(chart_info.get("title", "").upper(),
                  color=text_color, fontweight='bold', loc='left', pad=15)
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=150,
                    bbox_inches='tight', facecolor=bg_color)
        img_buffer.seek(0)
        plt.close(fig)
        return img_buffer
    except:
        return None


def add_image_placeholder(slide, prompt_text, theme_data):
    accent_rgb = hex_to_rgb(theme_data["accent_hex"])
    text_rgb = hex_to_rgb(theme_data["text_hex"])

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(
        7), Inches(2), Inches(5.5), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*accent_rgb)
    shape.fill.transparency = 0.8
    shape.line.color.rgb = RGBColor(*accent_rgb)

    tf = shape.text_frame
    tf.text = f"🖼️ AI VISUAL IDEA:\n\n{prompt_text}"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(*text_rgb)
        p.font.bold = True

# --- 5. PPTX BUILDER ---


def create_ppt_from_json(json_str, theme_data, google_key):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_rgb = hex_to_rgb(theme_data["bg_hex"])
    text_rgb = hex_to_rgb(theme_data["text_hex"])
    accent_rgb = hex_to_rgb(theme_data["accent_hex"])
    theme_name = theme_data.get("theme_name", "Modern")

    clean_json = json_str.replace("```json", "").replace("```", "").strip()
    try:
        slides_data = json.loads(clean_json)
    except:
        return None

    if 'image_prompt' in str(clean_json):
        prog_bar = st.progress(0, text="Generating Slides & Visuals...")
        total_slides = len(slides_data)
    else:
        prog_bar = None

    for i, slide_info in enumerate(slides_data):
        if prog_bar:
            prog_bar.progress((i + 1) / total_slides,
                              text=f"Processing Slide {i+1}...")

        is_section = slide_info.get("type") == "section"
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(
            *accent_rgb) if is_section else RGBColor(*bg_rgb)
        current_text_rgb = RGBColor(
            255, 255, 255) if is_section else RGBColor(*text_rgb)

        if is_section:
            tb = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11.3), Inches(2))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            apply_markdown_to_paragraph(
                p, slide_info.get("title", ""), current_text_rgb)
            for run in p.runs:
                run.font.size = Pt(54)
            p.alignment = PP_ALIGN.CENTER
        else:
            tb = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            apply_markdown_to_paragraph(
                p, slide_info.get("title", ""), current_text_rgb)
            for run in p.runs:
                run.font.size = Pt(32)

            has_chart = slide_info.get("chart") is not None
            has_image_prompt = slide_info.get("image_prompt") is not None
            width = Inches(6) if (
                has_chart or has_image_prompt) else Inches(11)

            cb = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), width, Inches(5))
            tf = cb.text_frame
            tf.word_wrap = True
            for point in slide_info.get("points", []):
                p = tf.add_paragraph()
                apply_markdown_to_paragraph(p, point, current_text_rgb)
                p.space_after = Pt(14)

            visual_placed = False
            if has_chart:
                img = create_chart_image(slide_info["chart"], theme_data)
                if img:
                    slide.shapes.add_picture(img, Inches(
                        7), Inches(2), width=Inches(5.5))
                    visual_placed = True
            if has_image_prompt and not visual_placed:
                ai_img = generate_real_image(
                    slide_info["image_prompt"], theme_name, google_key)
                if ai_img:
                    slide.shapes.add_picture(ai_img, Inches(
                        7), Inches(2), width=Inches(5.5))
                    visual_placed = True
                else:
                    add_image_placeholder(
                        slide, slide_info["image_prompt"], theme_data)

    if prog_bar:
        prog_bar.empty()
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer


# --- 6. INTERFACE ---
with st.sidebar:
    st.markdown("### 🛠️ Settings")
    if "GOOGLE_API_KEY" in st.secrets:
        google_key = st.secrets["GOOGLE_API_KEY"]
    else:
        google_key = st.text_input("Google API Key", type="password")

    st.divider()

    # NEW: Direct API Test (No Library Needed)
    if st.button("Test Image Gen (Direct API)"):
        test_url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-3.0-generate-001:predict?key={google_key}"
        try:
            r = requests.post(test_url, json={"instances": [
                              {"prompt": "Robot"}], "parameters": {"sampleCount": 1}})
            if r.status_code == 200:
                st.success("✅ Connection Successful!")
                b64 = r.json()['predictions'][0]['bytesBase64Encoded']
                st.image(base64.b64decode(b64), width=150)
            else:
                st.error(f"❌ API Error: {r.status_code}")
                st.caption(r.text)
        except Exception as e:
            st.error(f"❌ Connection Error: {e}")

    if st.button("Clear Cache"):
        st.cache_data.clear()
        st.rerun()

    st.divider()
    simulate_data = st.checkbox("Simulate Missing Data?", value=True)

st.title("AI Presentation Architect")
st.caption("Director Mode • Pure Google Ecosystem")

if not google_key:
    st.stop()

if "deck_json" not in st.session_state:
    st.session_state.deck_json = None
if "theme_data" not in st.session_state:
    st.session_state.theme_data = None
if "deck_topic" not in st.session_state:
    st.session_state.deck_topic = ""
if "ppt_binary" not in st.session_state:
    st.session_state.ppt_binary = None

with st.container(border=True):
    col1, col2 = st.columns([1, 1])
    with col1:
        st.subheader("1. Source Material")
        files = st.file_uploader(
            "Upload Docs/Data", accept_multiple_files=True)
        notes = st.text_area("Or Paste Text", height=100)
    with col2:
        st.subheader("2. Director's Notes")
        topic = st.text_input("Presentation Title",
                              placeholder="e.g. Q1 Boba Sales")
        audience = st.text_input(
            "Target Audience", placeholder="e.g. Investors")
        style_guide = st.text_area("Custom Instructions", height=100,
                                   placeholder="e.g. 'Professional tone. Minimal emojis.'")
        slide_count = st.slider("Target Slide Count", 3, 10, 6)

if st.button("Generate Deck", type="primary"):
    if (files or notes) and topic:
        with st.status("Architecting Presentation...", expanded=True) as status:
            st.write("🎨 Designing theme...")
            theme_data = generate_design_theme(google_key, topic, audience)
            st.session_state.theme_data = theme_data

            text = ""
            if files:
                for f in files:
                    text += extract_text_from_file(f)
            if notes:
                text += f"\n\nNOTES: {notes}"

            st.write("✍️ Writing content...")
            chart_rule = "GENERATE A CHART with estimated values if numbers are missing." if simulate_data else "No charts unless numbers exist."

            prompt = f"""
            Role: Expert Presentation Designer. Topic: {topic} | Audience: {audience}
            USER INSTRUCTIONS: "{style_guide}"
            TARGET LENGTH: EXACTLY {slide_count} SLIDES.
            SOURCE: {text[:30000]}
            
            RULES:
            1. Generate exactly {slide_count} slides.
            2. Use Markdown bold (**text**) for emphasis.
            3. {chart_rule}
            4. If no chart, add "image_prompt": "Description of visual..."
            5. **EMOJI RULE**: SPARINGLY. Only 1 per slide title if relevant.
            
            FORMAT (JSON ONLY):
            [ {{ 
                "type": "content", 
                "title": "Slide **Title**", 
                "points": ["Point 1", "Point 2"], 
                "chart": {{ "type": "BAR", "categories": ["A", "B"], "values": [10, 20], "title": "Sales" }},
                "image_prompt": "Visual description..."
            }} ]
            """
            response = get_gemini_response(google_key, prompt, temp=0.5)
            try:
                st.session_state.deck_json = json.loads(
                    response.replace("```json", "").replace("```", "").strip())
                st.session_state.deck_topic = topic

                st.write("🎨 Synthesizing Visuals...")
                st.session_state.ppt_binary = create_ppt_from_json(
                    response, theme_data, google_key)

                status.update(label="Complete",
                              state="complete", expanded=False)
            except Exception as e:
                st.error("AI Error")
                st.write(f"Error: {e}")
                st.expander("Raw Output").write(response)

if st.session_state.deck_json and st.session_state.theme_data:
    st.divider()
    td = st.session_state.theme_data
    st.subheader(f"🎞️ Preview: {st.session_state.deck_topic}")

    if st.session_state.ppt_binary:
        st.download_button("Download PowerPoint (.pptx)", st.session_state.ppt_binary,
                           f"{st.session_state.deck_topic}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary")

    bg_color = td["bg_hex"]
    text_color = td["text_hex"]
    accent_color = td["accent_hex"]

    for slide in st.session_state.deck_json:
        is_section = slide.get("type") == "section"
        has_chart = slide.get("chart") is not None
        has_image_prompt = slide.get("image_prompt") is not None

        card_bg = accent_color if is_section else bg_color
        card_text = "#FFFFFF" if is_section else text_color
        align = "center" if is_section else "left"

        with st.container():
            st.markdown(f"""<div style="background-color: {card_bg}; color: {card_text}; padding: 25px; border-radius: 10px; border: 1px solid #ddd; margin-bottom: 20px;">
                <h3 style="color: {card_text}; margin-top: 0; text-align: {align};">{markdown_to_html(slide.get('title', ''))}</h3>""", unsafe_allow_html=True)

            if not is_section:
                c1, c2 = st.columns([2, 1]) if (
                    has_chart or has_image_prompt) else st.columns([1, 0.01])
                with c1:
                    ul = "".join(
                        [f"<li style='font-size:18px; margin-bottom:8px;'>{markdown_to_html(p)}</li>" for p in slide.get("points", [])])
                    st.markdown(f"<ul>{ul}</ul>", unsafe_allow_html=True)
                with c2:
                    if has_chart:
                        img = create_chart_image(slide["chart"], td)
                        if img:
                            st.image(img, use_container_width=True)
                    elif has_image_prompt:
                        st.markdown(
                            f"<div style='text-align:center; padding:20px; border:2px dashed {accent_color}; color:{card_text}'>🖼️ AI Image (See Download)</div>", unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)
